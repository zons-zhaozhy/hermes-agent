"""End-to-end tests for the powerpoint skill helper scripts.

Runs entirely offline. Exercises create, read, template-fill (with
non-ASCII values), and edit (text + chart data + remove/move slide).
"""
import json
import os
import subprocess
import sys

import pytest

SKILL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCRIPTS = os.path.join(SKILL, "scripts")


def run(script, *args):
    env = dict(os.environ, LC_ALL="C", PYTHONIOENCODING="utf-8")
    proc = subprocess.run(
        [sys.executable, os.path.join(SCRIPTS, script), *args],
        capture_output=True, text=True, encoding="utf-8", env=env)
    assert proc.returncode == 0, f"{script} failed: {proc.stderr}"
    return json.loads(proc.stdout)


def write_json(path, obj):
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(obj, fh, ensure_ascii=False)


@pytest.fixture(scope="module")
def workdir(tmp_path_factory):
    return tmp_path_factory.mktemp("pptx")


@pytest.fixture(scope="module")
def sample_png(workdir):
    # 1x1 red PNG, hardcoded bytes -> no Pillow dependency.
    import base64
    data = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4"
        "z8DwHwAFBQIAX8jx0gAAAABJRU5ErkJggg==")
    path = workdir / "red.png"
    path.write_bytes(data)
    return str(path)


@pytest.fixture(scope="module")
def deck(workdir, sample_png):
    """Create a deck exercising every create feature."""
    spec = {
        "slide_size": "16:9",
        "slides": [
            {"layout": "title", "title": "Annual Review",
             "subtitle": "Fiscal year 2026",
             "notes": "Welcome the audience."},
            {"layout": "title_content", "title": "Agenda",
             "bullets": [
                 "Overview",
                 {"text": "Details", "level": 1, "bold": True,
                  "size": 20, "color": "CC0000", "font": "Arial"},
                 {"text": "Deep dive", "level": 2, "italic": True},
             ],
             "notes": "Keep this under two minutes."},
            {"layout": "blank", "title": None,
             "images": [{"path": sample_png, "left": 0.5, "top": 0.5,
                         "width": 2, "height": 2}],
             "tables": [{"left": 3, "top": 1, "width": 6, "height": 2,
                         "rows": [["Region", "Sales"],
                                  ["North", "120"], ["South", "80"]]}],
             "shapes": [{"type": "rounded_rectangle", "left": 10,
                         "top": 1, "width": 2.5, "height": 1,
                         "fill": "4472C4", "text": "Callout",
                         "text_color": "FFFFFF"}]},
            {"layout": "title_only", "title": "Charts",
             "charts": [
                 {"type": "bar", "left": 0.5, "top": 1.5, "width": 4,
                  "height": 3, "title": "Sales",
                  "categories": ["Q1", "Q2"],
                  "series": {"North": [10, 20], "South": [7, 13]}},
                 {"type": "line", "left": 4.7, "top": 1.5, "width": 4,
                  "height": 3, "categories": ["Jan", "Feb", "Mar"],
                  "series": {"Trend": [1, 3, 2]}},
                 {"type": "pie", "left": 8.9, "top": 1.5, "width": 4,
                  "height": 3, "categories": ["A", "B", "C"],
                  "series": {"Share": [50, 30, 20]}}]},
        ],
    }
    spec_path = workdir / "deck.json"
    write_json(spec_path, spec)
    out = workdir / "deck.pptx"
    result = run("pptx_create.py", str(spec_path), str(out))
    assert result["ok"] and result["slides"] == 4
    return str(out)


def test_create_and_outline(deck):
    outline = run("pptx_read.py", deck, "--outline")
    assert outline["slide_count"] == 4
    assert outline["slide_size_inches"][0] == pytest.approx(13.333, abs=0.01)
    s0, s1, s2, s3 = outline["slides"]
    assert "Annual Review" in s0["texts"]
    assert "Fiscal year 2026" in s0["texts"]
    assert s0["notes"] == "Welcome the audience."
    # bullets present on slide 1
    assert any("Deep dive" in t for t in s1["texts"])
    # table content
    assert s2["tables"][0][0] == ["Region", "Sales"]
    assert s2["tables"][0][2] == ["South", "80"]
    # image inventory + shape text
    assert len(s2["images"]) == 1 and s2["images"][0]["ext"] == "png"
    assert "Callout" in s2["texts"]
    # three charts with correct data
    assert len(s3["charts"]) == 3
    bar = s3["charts"][0]
    assert bar["categories"] == ["Q1", "Q2"]
    north = next(s for s in bar["series"] if s["name"] == "North")
    assert north["values"] == [10.0, 20.0]


def test_create_43_size(workdir):
    spec_path = workdir / "small.json"
    write_json(spec_path, {"slide_size": "4:3",
                           "slides": [{"layout": "title", "title": "T"}]})
    out = workdir / "small.pptx"
    run("pptx_create.py", str(spec_path), str(out))
    outline = run("pptx_read.py", str(out))
    assert outline["slide_size_inches"] == [10.0, 7.5]


def test_notes_mode(deck):
    result = run("pptx_read.py", deck, "--notes")
    assert result["notes"][1] == "Keep this under two minutes."


def test_image_export(deck, workdir):
    out_dir = workdir / "exported"
    result = run("pptx_read.py", deck, "--images", str(out_dir))
    assert len(result["exported"]) == 1
    exported = result["exported"][0]
    assert os.path.getsize(exported) > 0
    with open(exported, "rb") as fh:
        assert fh.read(8) == b"\x89PNG\r\n\x1a\n"


def test_template_fill_non_ascii(workdir):
    """Template token fill with non-ASCII values, forced under LC_ALL=C."""
    tpl_spec = workdir / "tpl.json"
    write_json(tpl_spec, {"slides": [
        {"layout": "title", "title": "{{city}} report",
         "subtitle": "Prepared by {{author}}",
         "notes": "Deck for {{city}}."},
        {"layout": "title_content", "title": "Data",
         "tables": [{"rows": [["Site", "{{city}}"]]}]},
    ]})
    template = workdir / "template.pptx"
    run("pptx_create.py", str(tpl_spec), str(template))

    values = workdir / "values.json"
    write_json(values, {"city": "Z\u00fcrich \u2014 \u2018Bericht\u2019",
                        "author": "Beispiel GmbH"})
    out = workdir / "filled.pptx"
    result = run("pptx_from_template.py", str(template), str(out),
                 "--values", str(values))
    assert result["tokens_filled"] == 4

    outline = run("pptx_read.py", str(out))
    expected = "Z\u00fcrich \u2014 \u2018Bericht\u2019"
    assert f"{expected} report" in outline["slides"][0]["texts"]
    assert outline["slides"][0]["notes"] == f"Deck for {expected}."
    assert outline["slides"][1]["tables"][0][0][1] == expected


def test_template_add_slides(workdir):
    template = workdir / "template.pptx"
    add_spec = workdir / "add.json"
    write_json(add_spec, {"slides": [
        {"layout": 1, "title": "Appended", "bullets": ["from layout"],
         "notes": "appended slide"}]})
    out = workdir / "appended.pptx"
    run("pptx_from_template.py", str(template), str(out),
        "--add-slides", str(add_spec))
    outline = run("pptx_read.py", str(out))
    assert outline["slide_count"] == 3
    assert "Appended" in outline["slides"][2]["texts"]


def test_edit_replace_text(deck, workdir):
    edited = workdir / "edited.pptx"
    result = run("pptx_edit.py", deck, "--output", str(edited),
                 "--replace-text", "Annual Review", "Semi-Annual Review",
                 "--replace-text", "South", "West")
    assert result["replacements"] == 2  # title + table cell
    outline = run("pptx_read.py", str(edited))
    assert "Semi-Annual Review" in outline["slides"][0]["texts"]
    assert outline["slides"][2]["tables"][0][2][0] == "West"
    # formatting survives a run-level replace: red bold bullet untouched
    assert any("Deep dive" in t for t in outline["slides"][1]["texts"])


def test_edit_chart_data(deck, workdir):
    spec = workdir / "chart_update.json"
    write_json(spec, {"slide": 3, "chart": 0,
                      "categories": ["Q3", "Q4"],
                      "series": {"North": [30, 40], "South": [21, 34]}})
    edited = workdir / "chart_edited.pptx"
    run("pptx_edit.py", deck, "--output", str(edited),
        "--chart-data", str(spec))
    outline = run("pptx_read.py", str(edited))
    bar = outline["slides"][3]["charts"][0]
    assert bar["categories"] == ["Q3", "Q4"]
    north = next(s for s in bar["series"] if s["name"] == "North")
    assert north["values"] == [30.0, 40.0]


def test_edit_remove_and_move_slide(deck, workdir):
    edited = workdir / "reordered.pptx"
    run("pptx_edit.py", deck, "--output", str(edited),
        "--remove-slide", "2", "--move-slide", "2", "0")
    outline = run("pptx_read.py", str(edited))
    assert outline["slide_count"] == 3
    # charts slide (was index 3, then 2 after removal) moved to front
    assert len(outline["slides"][0]["charts"]) == 3
    assert "Annual Review" in outline["slides"][1]["texts"]


def test_edit_swap_image(deck, workdir):
    import base64
    blue = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGNg"
        "YPj/HwADAgH/p5UronAAAAAASUVORK5CYII=")
    blue_path = workdir / "blue.png"
    blue_path.write_bytes(blue)
    outline = run("pptx_read.py", deck)
    # find picture shape name via python-pptx directly
    sys.path.insert(0, SCRIPTS)
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(deck)
    name = next(s.name for s in prs.slides[2].shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    edited = workdir / "swapped.pptx"
    run("pptx_edit.py", deck, "--output", str(edited),
        "--swap-image", "2", name, str(blue_path))
    out_dir = workdir / "swapped_images"
    result = run("pptx_read.py", str(edited), "--images", str(out_dir))
    with open(result["exported"][0], "rb") as fh:
        assert fh.read() == blue


def test_help_flags():
    for script in ("pptx_create.py", "pptx_read.py", "pptx_edit.py",
                   "pptx_from_template.py"):
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, script), "--help"],
            capture_output=True, text=True, encoding="utf-8")
        assert proc.returncode == 0 and "usage" in proc.stdout.lower()


# ---------------------------------------------------------------------------
# parity features: render, run-merge replace, chart ops, duplication,
# polish (background/hyperlink/slide-number/footer), notes editing
# ---------------------------------------------------------------------------

def run_raw(script, *args):
    env = dict(os.environ, LC_ALL="C", PYTHONIOENCODING="utf-8")
    return subprocess.run(
        [sys.executable, os.path.join(SCRIPTS, script), *args],
        capture_output=True, text=True, encoding="utf-8", env=env)


def test_render_all_slides(workdir):
    import shutil
    spec_path = workdir / "render_spec.json"
    write_json(spec_path, {"slides": [
        {"layout": "title", "title": "One"},
        {"layout": "title_content", "title": "Two", "bullets": ["b"]},
        {"layout": "blank"}]})
    deck3 = workdir / "render_me.pptx"
    run("pptx_create.py", str(spec_path), str(deck3))
    out_dir = workdir / "render_out"
    result = run("pptx_render.py", str(deck3), "--outdir", str(out_dir))
    have_tools = shutil.which("soffice") and (
        shutil.which("pdftoppm") or shutil.which("pdftocairo"))
    if have_tools:
        assert result["rendered"] is True
        assert len(result["files"]) == 3
        for png in result["files"]:
            with open(png, "rb") as fh:
                assert fh.read(8) == b"\x89PNG\r\n\x1a\n"
    else:
        assert result["rendered"] is False
        assert result["missing"]
        assert "guidance" in result


def test_replace_across_identically_formatted_runs(workdir):
    """A match split mid-word into equal-format runs keeps formatting."""
    import copy as cp
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1),
                                   Inches(6), Inches(1))
    para = box.text_frame.paragraphs[0]
    run_obj = para.add_run()
    run_obj.text = "Say TotalWord now"
    run_obj.font.bold = True
    run_obj.font.size = Pt(20)
    run_obj.font.color.rgb = RGBColor.from_string("CC0000")
    # simulate PowerPoint's spell-check split: same rPr, seam mid-match
    second = cp.deepcopy(run_obj._r)
    run_obj._r.addnext(second)
    run_obj.text = "Say Total"
    para.runs[1].text = "Word now"
    path = workdir / "split_runs.pptx"
    prs.save(str(path))

    edited = workdir / "split_runs_edited.pptx"
    result = run("pptx_edit.py", str(path), "--output", str(edited),
                 "--replace-text", "TotalWord", "MergedWord")
    assert result["replacements"] == 1

    prs2 = Presentation(str(edited))
    para2 = prs2.slides[0].shapes[0].text_frame.paragraphs[0]
    assert "".join(r.text for r in para2.runs) == "Say MergedWord now"
    for r in para2.runs:
        assert r.font.bold is True
        assert r.font.size == Pt(20)
        assert str(r.font.color.rgb) == "CC0000"


def test_chart_surgical_ops(deck, workdir):
    from pptx import Presentation
    spec = workdir / "chart_ops.json"
    write_json(spec, {"slide": 3, "chart": 0, "ops": [
        {"op": "update_series", "name": "North", "values": [99, 88]},
        {"op": "add_series", "name": "East", "values": [1, 2]},
        {"op": "remove_series", "name": "South"},
        {"op": "rename_category", "from": "Q1", "to": "Q1 FY26"},
        {"op": "set_title", "title": "Updated Sales"}]})
    edited = workdir / "chart_ops.pptx"
    run("pptx_edit.py", deck, "--output", str(edited),
        "--chart-data", str(spec))
    outline = run("pptx_read.py", str(edited))
    bar = outline["slides"][3]["charts"][0]
    assert bar["categories"] == ["Q1 FY26", "Q2"]
    names = [s["name"] for s in bar["series"]]
    assert "South" not in names and "East" in names
    north = next(s for s in bar["series"] if s["name"] == "North")
    assert north["values"] == [99.0, 88.0]
    east = next(s for s in bar["series"] if s["name"] == "East")
    assert east["values"] == [1.0, 2.0]
    chart = next(s.chart for s in Presentation(str(edited)).slides[3].shapes
                 if s.has_chart)
    assert chart.chart_title.text_frame.text == "Updated Sales"


def test_chart_op_unknown_series_fails(deck, workdir):
    spec = workdir / "chart_bad.json"
    write_json(spec, {"slide": 3, "chart": 0, "ops": [
        {"op": "update_series", "name": "Nowhere", "values": [1, 2]}]})
    proc = run_raw("pptx_edit.py", deck, "--output",
                   str(workdir / "never.pptx"), "--chart-data", str(spec))
    assert proc.returncode != 0
    assert "Nowhere" in proc.stderr


def test_duplicate_image_slide_is_independent(deck, workdir):
    import base64
    dup = workdir / "dup.pptx"
    result = run("pptx_edit.py", deck, "--output", str(dup),
                 "--duplicate-slide", "2")
    assert result["duplicated_to"] == 4
    outline = run("pptx_read.py", str(dup))  # re-opens cleanly
    assert outline["slide_count"] == 5
    s4 = outline["slides"][4]
    assert len(s4["images"]) == 1
    assert s4["tables"][0][0] == ["Region", "Sales"]
    assert "Callout" in s4["texts"]

    # independence: swap the image on the COPY, original stays red
    blue = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGNg"
        "YPj/HwADAgH/p5UronAAAAAASUVORK5CYII=")
    blue_path = workdir / "blue2.png"
    blue_path.write_bytes(blue)
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(dup))
    name = next(s.name for s in prs.slides[4].shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    run("pptx_edit.py", str(dup), "--swap-image", "4", name,
        str(blue_path))
    prs = Presentation(str(dup))
    orig = next(s for s in prs.slides[2].shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    copy_pic = next(s for s in prs.slides[4].shapes
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert copy_pic.image.blob == blue
    assert orig.image.blob != blue


def test_duplicate_chart_slide_refused(deck, workdir):
    proc = run_raw("pptx_edit.py", deck, "--output",
                   str(workdir / "never2.pptx"), "--duplicate-slide", "3")
    assert proc.returncode != 0
    assert "chart" in proc.stderr.lower()


def test_create_polish_features(workdir):
    from pptx import Presentation
    spec_path = workdir / "polish.json"
    write_json(spec_path, {"slides": [
        {"layout": "title_content", "title": "Polished",
         "background": "112233", "footer": "Confidential draft",
         "slide_number": True,
         "bullets": [{"text": "Visit example",
                      "link": "https://example.com/info"}]}]})
    out = workdir / "polish.pptx"
    run("pptx_create.py", str(spec_path), str(out))
    prs = Presentation(str(out))
    slide = prs.slides[0]
    assert str(slide.background.fill.fore_color.rgb) == "112233"
    ph_idx = [ph.placeholder_format.idx for ph in slide.placeholders]
    assert 12 in ph_idx  # slide number enabled
    footer = next(ph for ph in slide.placeholders
                  if ph.placeholder_format.idx == 11)
    assert footer.text_frame.text == "Confidential draft"
    link_runs = [r for sh in slide.shapes if sh.has_text_frame
                 for p in sh.text_frame.paragraphs for r in p.runs
                 if r.hyperlink.address]
    assert link_runs[0].hyperlink.address == "https://example.com/info"


def test_edit_polish_features(deck, workdir):
    from pptx import Presentation
    edited = workdir / "polish_edit.pptx"
    result = run("pptx_edit.py", deck, "--output", str(edited),
                 "--set-background", "0", "004400",
                 "--hyperlink", "1", "Overview", "https://example.com/x",
                 "--enable-slide-number", "1",
                 "--set-footer", "1", "Footer via edit")
    assert result["background_set"] and result["footer_set"]
    assert result["hyperlinked_runs"] == 1
    assert result["slide_number_enabled"] is True
    prs = Presentation(str(edited))
    assert str(prs.slides[0].background.fill.fore_color.rgb) == "004400"
    s1 = prs.slides[1]
    assert any(ph.placeholder_format.idx == 12 for ph in s1.placeholders)
    footer = next(ph for ph in s1.placeholders
                  if ph.placeholder_format.idx == 11)
    assert footer.text_frame.text == "Footer via edit"
    links = [r.hyperlink.address for sh in s1.shapes if sh.has_text_frame
             for p in sh.text_frame.paragraphs for r in p.runs
             if r.hyperlink.address]
    assert links == ["https://example.com/x"]


def test_set_and_append_notes(deck, workdir):
    edited = workdir / "notes_edit.pptx"
    result = run("pptx_edit.py", deck, "--output", str(edited),
                 "--set-notes", "0", "Fresh notes")
    assert result["notes_set"]
    run("pptx_edit.py", str(edited), "--append-notes", "0", "Second line")
    notes = run("pptx_read.py", str(edited), "--notes")
    assert notes["notes"][0] == "Fresh notes\nSecond line"


def test_render_help_flag():
    proc = run_raw("pptx_render.py", "--help")
    assert proc.returncode == 0 and "usage" in proc.stdout.lower()
