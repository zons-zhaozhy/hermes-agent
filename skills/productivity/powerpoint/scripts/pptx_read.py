#!/usr/bin/env python3
"""Read a .pptx file: JSON outline, notes, or export embedded images.

Modes:
  --outline      JSON with per-slide layout, texts, tables, notes,
                 chart data, and image inventory (default mode).
  --notes        JSON list of speaker notes per slide.
  --images DIR   Export every embedded picture to DIR as files.
"""
import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def iter_shapes(shapes):
    """Yield shapes, descending into groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def chart_info(chart):
    info = {"type": str(chart.chart_type),
            "categories": [str(c) for c in chart.plots[0].categories],
            "series": []}
    for plot in chart.plots:
        for series in plot.series:
            try:
                name = series.name
            except (AttributeError, KeyError):
                name = None
            info["series"].append({"name": name,
                                   "values": list(series.values)})
    return info


def slide_record(index, slide):
    rec = {"index": index, "layout": slide.slide_layout.name,
           "texts": [], "tables": [], "images": [], "charts": [],
           "notes": None}
    for shape in iter_shapes(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip():
            rec["texts"].append(shape.text_frame.text)
        if shape.has_table:
            rec["tables"].append(
                [[cell.text for cell in row.cells]
                 for row in shape.table.rows])
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                rec["images"].append({"filename": img.filename,
                                      "ext": img.ext,
                                      "size_bytes": len(img.blob)})
            except (KeyError, ValueError):
                rec["images"].append({"filename": None, "ext": None,
                                      "size_bytes": None,
                                      "note": "linked or unreadable"})
        if shape.has_chart:
            rec["charts"].append(chart_info(shape.chart))
    if slide.has_notes_slide:
        rec["notes"] = slide.notes_slide.notes_text_frame.text
    return rec


def export_images(prs, out_dir):
    os.makedirs(out_dir, exist_ok=True)
    written = []
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(iter_shapes(slide.shapes)):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img = shape.image
            except (KeyError, ValueError):
                continue
            path = os.path.join(out_dir, f"slide{i}_img{j}.{img.ext}")
            with open(path, "wb") as fh:
                fh.write(img.blob)
            written.append(path)
    return written


def main(argv=None):
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(
        description="Read a .pptx: outline/notes as JSON, export images.")
    parser.add_argument("pptx", help="path to the .pptx file")
    parser.add_argument("--outline", action="store_true",
                        help="print full JSON outline (default)")
    parser.add_argument("--notes", action="store_true",
                        help="print speaker notes only")
    parser.add_argument("--images", metavar="DIR",
                        help="export embedded images into DIR")
    args = parser.parse_args(argv)

    prs = Presentation(args.pptx)

    if args.images:
        written = export_images(prs, args.images)
        print(json.dumps({"ok": True, "exported": written}, indent=2))
        return 0
    if args.notes:
        notes = [slide.notes_slide.notes_text_frame.text
                 if slide.has_notes_slide else None
                 for slide in prs.slides]
        print(json.dumps({"ok": True, "notes": notes},
                         indent=2, ensure_ascii=True))
        return 0

    outline = {
        "ok": True,
        "slide_size_inches": [round(Emu(prs.slide_width).inches, 3),
                              round(Emu(prs.slide_height).inches, 3)],
        "slide_count": len(prs.slides._sldIdLst),
        "layouts_available": [lay.name for lay in prs.slide_layouts],
        "slides": [slide_record(i, s) for i, s in enumerate(prs.slides)],
    }
    print(json.dumps(outline, indent=2, ensure_ascii=True))
    return 0


if __name__ == "__main__":
    sys.exit(main())
