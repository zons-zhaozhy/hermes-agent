#!/usr/bin/env python3
"""Create a .pptx presentation from a JSON deck spec.

Spec format (all positions/sizes in inches, colors as RRGGBB hex):
{
  "slide_size": "16:9",              // or "4:3" (default "16:9")
  "slides": [
    {"layout": "title", "title": "My Deck", "subtitle": "Q3 review"},
    {"layout": "title_content", "title": "Agenda",
     "bullets": ["Top item",
                 {"text": "Sub item", "level": 1, "bold": true,
                  "size": 18, "color": "CC0000", "font": "Arial",
                  "italic": false,
                  "link": "https://example.com/agenda"}],
     "background": "1F2937",           // solid slide background (hex)
     "footer": "Confidential",         // footer placeholder text
     "slide_number": true,             // enable slide-number placeholder
     "notes": "Speaker notes for this slide"},
    {"layout": "blank", "title": "Widgets",
     "images":  [{"path": "logo.png", "left": 1, "top": 1, "width": 3}],
     "tables":  [{"left": 1, "top": 2, "width": 6, "height": 2,
                  "rows": [["H1", "H2"], ["a", "b"]]}],
     "shapes":  [{"type": "rounded_rectangle", "left": 8, "top": 1,
                  "width": 3, "height": 1, "fill": "4472C4",
                  "text": "Callout", "text_color": "FFFFFF"}],
     "charts":  [{"type": "bar", "left": 1, "top": 3, "width": 6,
                  "height": 3.5, "title": "Sales",
                  "categories": ["Q1", "Q2"],
                  "series": {"North": [10, 20], "South": [7, 13]}}]}
  ]
}
Layouts: title, title_content, section, two_content, title_only, blank
Chart types: bar, bar_h, line, pie   Shape types: rectangle,
rounded_rectangle, oval, diamond, right_arrow, chevron
"""
import argparse
import copy
import json
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

LAYOUTS = {"title": 0, "title_content": 1, "section": 2,
           "two_content": 3, "title_only": 5, "blank": 6}
CHART_TYPES = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
               "bar_h": XL_CHART_TYPE.BAR_CLUSTERED,
               "line": XL_CHART_TYPE.LINE_MARKERS,
               "pie": XL_CHART_TYPE.PIE}
SHAPE_TYPES = {"rectangle": MSO_SHAPE.RECTANGLE,
               "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
               "oval": MSO_SHAPE.OVAL, "diamond": MSO_SHAPE.DIAMOND,
               "right_arrow": MSO_SHAPE.RIGHT_ARROW,
               "chevron": MSO_SHAPE.CHEVRON}


def style_run(run, spec):
    """Apply font styling from a bullet/text spec dict to a run."""
    font = run.font
    if spec.get("size"):
        font.size = Pt(spec["size"])
    if spec.get("bold") is not None:
        font.bold = spec["bold"]
    if spec.get("italic") is not None:
        font.italic = spec["italic"]
    if spec.get("font"):
        font.name = spec["font"]
    if spec.get("color"):
        font.color.rgb = RGBColor.from_string(spec["color"])
    if spec.get("link"):
        run.hyperlink.address = spec["link"]


def add_bullets(text_frame, bullets):
    text_frame.clear()
    for i, item in enumerate(bullets):
        if isinstance(item, str):
            item = {"text": item}
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.level = int(item.get("level", 0))
        run = para.add_run()
        run.text = item.get("text", "")
        style_run(run, item)


def copy_layout_placeholder(slide, ph_idx):
    """Copy a layout placeholder (footer=11, slide number=12) onto the
    slide so it actually renders; returns the shape or None if the layout
    does not provide it."""
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            slide.shapes._spTree.append(copy.deepcopy(ph._element))
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == ph_idx:
                    return shape
    return None


def build_slide(prs, spec):
    layout_idx = LAYOUTS.get(spec.get("layout", "title_content"), 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    if spec.get("background"):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(spec["background"])
    if spec.get("slide_number"):
        copy_layout_placeholder(slide, 12)
    if spec.get("footer"):
        shape = copy_layout_placeholder(slide, 11)
        if shape is not None:
            shape.text_frame.text = spec["footer"]

    if spec.get("title") is not None and slide.shapes.title is not None:
        slide.shapes.title.text = spec["title"]
    if spec.get("subtitle") is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = spec["subtitle"]
                break
    if spec.get("bullets"):
        body = next((ph for ph in slide.placeholders
                     if ph.placeholder_format.idx != 0), None)
        if body is None:
            body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5),
                                            Inches(9), Inches(5))
        add_bullets(body.text_frame, spec["bullets"])

    for img in spec.get("images", []):
        kwargs = {}
        if img.get("width"):
            kwargs["width"] = Inches(img["width"])
        if img.get("height"):
            kwargs["height"] = Inches(img["height"])
        slide.shapes.add_picture(img["path"], Inches(img.get("left", 1)),
                                 Inches(img.get("top", 1)), **kwargs)

    for tbl in spec.get("tables", []):
        rows = tbl["rows"]
        shape = slide.shapes.add_table(
            len(rows), len(rows[0]), Inches(tbl.get("left", 1)),
            Inches(tbl.get("top", 2)), Inches(tbl.get("width", 6)),
            Inches(tbl.get("height", 2)))
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                shape.table.cell(r, c).text = str(val)

    for shp in spec.get("shapes", []):
        shape = slide.shapes.add_shape(
            SHAPE_TYPES.get(shp.get("type", "rectangle"), MSO_SHAPE.RECTANGLE),
            Inches(shp.get("left", 1)), Inches(shp.get("top", 1)),
            Inches(shp.get("width", 2)), Inches(shp.get("height", 1)))
        if shp.get("fill"):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(shp["fill"])
        if shp.get("text"):
            shape.text_frame.text = shp["text"]
            if shp.get("text_color"):
                run = shape.text_frame.paragraphs[0].runs[0]
                run.font.color.rgb = RGBColor.from_string(shp["text_color"])

    for cht in spec.get("charts", []):
        data = CategoryChartData()
        data.categories = cht["categories"]
        for name, values in cht["series"].items():
            data.add_series(name, values)
        frame = slide.shapes.add_chart(
            CHART_TYPES.get(cht.get("type", "bar"),
                            XL_CHART_TYPE.COLUMN_CLUSTERED),
            Inches(cht.get("left", 1)), Inches(cht.get("top", 2)),
            Inches(cht.get("width", 6)), Inches(cht.get("height", 4)), data)
        if cht.get("title"):
            frame.chart.has_title = True
            frame.chart.chart_title.text_frame.text = cht["title"]

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    return slide


def main(argv=None):
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(
        description="Create a .pptx deck from a JSON spec.",
        epilog=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("spec", help="path to JSON deck spec")
    parser.add_argument("output", help="output .pptx path")
    args = parser.parse_args(argv)

    with open(args.spec, encoding="utf-8") as fh:
        spec = json.load(fh)

    prs = Presentation()
    if spec.get("slide_size", "16:9") == "16:9":
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    else:
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)

    for slide_spec in spec.get("slides", []):
        build_slide(prs, slide_spec)

    prs.save(args.output)
    print(json.dumps({"ok": True, "output": args.output,
                      "slides": len(prs.slides._sldIdLst)}))
    return 0


if __name__ == "__main__":
    sys.exit(main())
