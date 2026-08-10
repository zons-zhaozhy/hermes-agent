#!/usr/bin/env python3
"""Build a deck from a .pptx template (brand deck) and fill placeholders.

Two modes:
  1) Token fill (default): open TEMPLATE, replace every {{token}} across
     slides, tables, and notes using --values JSON ({"token": "value"}),
     save to OUTPUT. Formatting of the token's run is preserved.
  2) --add-slides SPEC.json: additionally append slides built from the
     template's own layouts (referenced by layout name or index), so new
     slides inherit the brand master. Spec:
     {"slides": [{"layout": "Title and Content", "title": "New",
                  "bullets": ["a", {"text": "b", "level": 1}],
                  "notes": "presenter text"}]}
"""
import argparse
import json
import sys

from pptx import Presentation


def fill_tokens(prs, values):
    from pptx_edit import replace_text  # same scripts/ directory
    total = 0
    for token, value in values.items():
        total += replace_text(prs, "{{%s}}" % token, str(value))
    return total


def find_layout(prs, ref):
    if isinstance(ref, int):
        return prs.slide_layouts[ref]
    for layout in prs.slide_layouts:
        if layout.name == ref:
            return layout
    raise SystemExit(f"layout {ref!r} not found; available: "
                     f"{[la.name for la in prs.slide_layouts]}")


def add_slides(prs, spec):
    from pptx_create import add_bullets
    for slide_spec in spec.get("slides", []):
        layout = find_layout(prs, slide_spec.get("layout", 1))
        slide = prs.slides.add_slide(layout)
        if slide_spec.get("title") is not None and slide.shapes.title:
            slide.shapes.title.text = slide_spec["title"]
        if slide_spec.get("bullets"):
            body = next((ph for ph in slide.placeholders
                         if ph.placeholder_format.idx != 0), None)
            if body is not None:
                add_bullets(body.text_frame, slide_spec["bullets"])
        if slide_spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_spec["notes"]


def main(argv=None):
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(
        description="Fill {{tokens}} in a .pptx template and optionally "
                    "append slides using the template's own layouts.",
        epilog=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("template", help="path to the template .pptx")
    parser.add_argument("output", help="output .pptx path")
    parser.add_argument("--values", metavar="JSON",
                        help="JSON file mapping token -> replacement value")
    parser.add_argument("--add-slides", metavar="SPEC_JSON",
                        help="JSON spec of slides to append")
    args = parser.parse_args(argv)

    prs = Presentation(args.template)
    filled = 0
    if args.values:
        with open(args.values, encoding="utf-8") as fh:
            filled = fill_tokens(prs, json.load(fh))
    if args.add_slides:
        with open(args.add_slides, encoding="utf-8") as fh:
            add_slides(prs, json.load(fh))
    prs.save(args.output)
    print(json.dumps({"ok": True, "output": args.output,
                      "tokens_filled": filled}))
    return 0


if __name__ == "__main__":
    import os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    sys.exit(main())
