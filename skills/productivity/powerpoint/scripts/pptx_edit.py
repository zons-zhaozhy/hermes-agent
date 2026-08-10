#!/usr/bin/env python3
"""Edit a .pptx in place (or save to --output).

Operations (repeatable / combinable):
  --replace-text OLD NEW   Replace text everywhere (slides, tables, notes).
                           Adjacent runs with identical formatting are
                           merged first, so matches PowerPoint split across
                           identically-formatted runs keep their formatting.
                           Only a match spanning genuinely different
                           formats falls back to a paragraph rewrite with
                           the first run's font (documented caveat).
  --chart-data SPEC.json   Update a chart. Full replace spec:
                           {"slide": 0, "chart": 0,
                            "categories": ["Q1", "Q2"],
                            "series": {"North": [1, 2], "South": [3, 4]}}
                           Or surgical ops (existing data is read, modified,
                           and written back via replace_data):
                           {"slide": 0, "chart": 0, "ops": [
                             {"op": "update_series", "name": "North",
                              "values": [5, 6]},
                             {"op": "add_series", "name": "East",
                              "values": [1, 2]},
                             {"op": "remove_series", "name": "South"},
                             {"op": "rename_category", "from": "Q1",
                              "to": "Q1 FY26"},
                             {"op": "set_title", "title": "New title"}]}
  --swap-image SLIDE SHAPE_NAME NEW_IMAGE
                           Replace a picture's bits, keeping position/size.
  --remove-slide N         Delete slide at index N (0-based).
  --move-slide FROM TO     Reorder: move slide FROM to position TO.
  --duplicate-slide N      Append an independent deep copy of slide N
                           (text, images, tables, shapes, notes). Refuses
                           slides containing charts (a chart embeds an XLSX
                           workbook part that cannot be cloned reliably).
  --set-background N HEX   Solid background color for slide N.
  --hyperlink N TEXT URL   Make runs containing TEXT on slide N links.
  --enable-slide-number N  Copy the layout's slide-number placeholder in.
  --set-footer N TEXT      Enable the layout's footer placeholder with TEXT.
  --set-notes N TEXT       Replace slide N's speaker notes.
  --append-notes N TEXT    Append a paragraph to slide N's speaker notes.
"""
import argparse
import copy
import json
import sys

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

R_EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _run_format_key(r):
    """Canonical string for a run's <a:rPr>; None when absent."""
    rPr = r.find(qn("a:rPr"))
    if rPr is None:
        return None
    return etree.tostring(rPr)


def normalize_runs(para):
    """Merge adjacent runs whose formatting is byte-identical.

    PowerPoint splits paragraph text into runs at spell-check and edit
    boundaries even when formatting never changes; merging them back makes
    cross-run text replacement lossless for the common case.
    """
    runs = list(para.runs)
    i = 0
    while i + 1 < len(runs):
        a, b = runs[i], runs[i + 1]
        if (_run_format_key(a._r) == _run_format_key(b._r)
                and a._r.getnext() is b._r):
            a.text = a.text + b.text
            b._r.getparent().remove(b._r)
            runs.pop(i + 1)
        else:
            i += 1


def replace_in_text_frame(text_frame, old, new):
    count = 0
    for para in text_frame.paragraphs:
        joined = "".join(run.text for run in para.runs)
        if old not in joined:
            continue
        if not any(old in run.text for run in para.runs):
            # Match spans runs: merge identically-formatted neighbours
            # first, which resolves pure spell-check splits losslessly.
            normalize_runs(para)
        if any(old in run.text for run in para.runs):
            # Run-level replace: preserves each run's formatting exactly.
            for run in para.runs:
                if old in run.text:
                    count += run.text.count(old)
                    run.text = run.text.replace(old, new)
        else:
            # Match spans genuinely differently-formatted runs -> rewrite
            # paragraph, keeping only the first run's formatting (caveat).
            joined = "".join(run.text for run in para.runs)
            count += joined.count(old)
            first = para.runs[0]
            first.text = joined.replace(old, new)
            for run in para.runs[1:]:
                run._r.getparent().remove(run._r)
    return count


def iter_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
    if slide.has_notes_slide:
        yield slide.notes_slide.notes_text_frame


def replace_text(prs, old, new):
    total = 0
    for slide in prs.slides:
        for tf in iter_text_frames(slide):
            total += replace_in_text_frame(tf, old, new)
    return total


def _read_chart_data(chart):
    """Current categories and ordered (name, values) pairs of a chart."""
    categories = [str(c) for c in chart.plots[0].categories]
    series = []
    for plot in chart.plots:
        for s in plot.series:
            try:
                name = s.name
            except (AttributeError, KeyError):
                name = ""
            series.append([name, list(s.values)])
    return categories, series


def update_chart(prs, spec_path):
    """Full replace ("categories"+"series") or surgical "ops".

    python-pptx can only swap a chart's entire dataset (replace_data), so
    surgical ops are implemented as read-existing -> modify -> replace.
    """
    with open(spec_path, encoding="utf-8") as fh:
        spec = json.load(fh)
    slide = prs.slides[spec.get("slide", 0)]
    charts = [s.chart for s in slide.shapes if s.has_chart]
    if not charts:
        raise SystemExit(f"no chart on slide {spec.get('slide', 0)}")
    chart = charts[spec.get("chart", 0)]

    if "ops" in spec:
        categories, series = _read_chart_data(chart)
        dirty = False
        for op in spec["ops"]:
            kind = op["op"]
            if kind == "update_series":
                match = [s for s in series if s[0] == op["name"]]
                if not match:
                    raise SystemExit(f"no series named {op['name']!r}")
                match[0][1] = op["values"]
                dirty = True
            elif kind == "add_series":
                series.append([op["name"], op["values"]])
                dirty = True
            elif kind == "remove_series":
                before = len(series)
                series = [s for s in series if s[0] != op["name"]]
                if len(series) == before:
                    raise SystemExit(f"no series named {op['name']!r}")
                dirty = True
            elif kind == "rename_category":
                if "index" in op:
                    idx = int(op["index"])
                else:
                    if op["from"] not in categories:
                        raise SystemExit(
                            f"no category named {op['from']!r}")
                    idx = categories.index(op["from"])
                categories[idx] = op["to"]
                dirty = True
            elif kind == "set_title":
                chart.has_title = True
                chart.chart_title.text_frame.text = op["title"]
            else:
                raise SystemExit(f"unknown chart op {kind!r}")
        if dirty:
            data = CategoryChartData()
            data.categories = categories
            for name, values in series:
                data.add_series(name, values)
            chart.replace_data(data)
        return

    data = CategoryChartData()
    data.categories = spec["categories"]
    for name, values in spec["series"].items():
        data.add_series(name, values)
    chart.replace_data(data)


def swap_image(prs, slide_idx, shape_name, new_path):
    slide = prs.slides[int(slide_idx)]
    for shape in slide.shapes:
        if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and shape.name == shape_name):
            image_part, rid = slide.part.get_or_add_image_part(new_path)
            blip = shape._element.blipFill.blip
            blip.set(R_EMBED + "embed", rid)
            return True
    raise SystemExit(f"no picture named {shape_name!r} on slide {slide_idx}")


def remove_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slide_id = list(sldIdLst)[int(index)]
    rid = slide_id.get(R_EMBED + "id")
    prs.part.drop_rel(rid)
    sldIdLst.remove(slide_id)


def move_slide(prs, src, dst):
    """Reorder by moving the <p:sldId> element inside <p:sldIdLst>."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    element = ids[int(src)]
    sldIdLst.remove(element)
    sldIdLst.insert(int(dst), element)


def duplicate_slide(prs, index):
    """Append an independent deep copy of slide `index`.

    Copies the shape tree XML and re-creates image/media relationships on
    the new slide part, remapping rIds. Charts are refused: each chart
    relationship embeds a separate XLSX workbook part, and cloning that
    graph reliably is not supported — better to refuse than corrupt.
    """
    source = prs.slides[int(index)]
    if any(sh.has_chart for sh in source.shapes):
        raise SystemExit(
            f"slide {index} contains a chart; duplication of chart slides "
            "is not supported (chart XML embeds a workbook part that "
            "cannot be cloned safely). Rebuild the chart on a new slide "
            "with pptx_create.py / pptx_from_template.py instead.")

    dest = prs.slides.add_slide(source.slide_layout)
    # drop the placeholders add_slide seeded from the layout
    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shape._element))

    # re-create the source slide's part relationships on the copy
    rid_map = {}
    for rel in list(source.part.rels.values()):
        if rel.reltype.endswith(("/slideLayout", "/notesSlide")):
            continue
        if rel.is_external:
            new_rid = dest.part.rels.get_or_add_ext_rel(
                rel.reltype, rel.target_ref)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel.rId] = new_rid

    for el in dest.shapes._spTree.iter():
        for attr, val in el.attrib.items():
            if attr.startswith(R_EMBED) and val in rid_map:
                el.set(attr, rid_map[val])

    if source.has_notes_slide:
        dest.notes_slide.notes_text_frame.text = (
            source.notes_slide.notes_text_frame.text)
    return len(prs.slides._sldIdLst) - 1


def set_background(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def add_hyperlink(prs, slide_idx, text, url):
    """Turn every run containing `text` on the slide into a hyperlink.

    The link applies to the whole run (python-pptx links whole runs).
    """
    slide = prs.slides[int(slide_idx)]
    hits = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if text in run.text:
                    run.hyperlink.address = url
                    hits += 1
    if not hits:
        raise SystemExit(f"no run containing {text!r} on slide {slide_idx}")
    return hits


def _copy_layout_placeholder(slide, ph_idx):
    """Copy the layout placeholder with idx `ph_idx` onto the slide.

    Slide-number (idx 12) and footer (idx 11) placeholders exist on the
    layout but are not inherited by a slide until the slide carries its
    own copy — this enables them. Returns the new shape, or None when the
    layout does not provide that placeholder.
    """
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            el = copy.deepcopy(ph._element)
            slide.shapes._spTree.append(el)
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == ph_idx:
                    return shape
            return None
    return None


def enable_slide_number(slide):
    if any(ph.placeholder_format.idx == 12 for ph in slide.placeholders):
        return True
    return _copy_layout_placeholder(slide, 12) is not None


def set_footer(slide, text):
    shape = next((ph for ph in slide.placeholders
                  if ph.placeholder_format.idx == 11), None)
    if shape is None:
        shape = _copy_layout_placeholder(slide, 11)
    if shape is None:
        raise SystemExit("layout provides no footer placeholder; add a "
                         "textbox instead")
    shape.text_frame.text = text
    return True


def set_notes(slide, text, append=False):
    tf = slide.notes_slide.notes_text_frame
    if append and tf.text:
        para = tf.add_paragraph()
        para.text = text
    else:
        tf.text = text


def main(argv=None):
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(
        description="Edit a .pptx: replace text, update chart data, swap "
                    "images, duplicate/remove/reorder slides, backgrounds, "
                    "hyperlinks, footers, slide numbers, speaker notes.",
        epilog=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", help="path to the .pptx file")
    parser.add_argument("--output", help="save to this path instead of "
                                         "overwriting the input")
    parser.add_argument("--replace-text", nargs=2, action="append",
                        metavar=("OLD", "NEW"), default=[])
    parser.add_argument("--chart-data", metavar="SPEC_JSON")
    parser.add_argument("--swap-image", nargs=3,
                        metavar=("SLIDE", "SHAPE_NAME", "IMAGE"))
    parser.add_argument("--remove-slide", type=int, metavar="N")
    parser.add_argument("--move-slide", nargs=2, type=int,
                        metavar=("FROM", "TO"))
    parser.add_argument("--duplicate-slide", type=int, metavar="N")
    parser.add_argument("--set-background", nargs=2,
                        metavar=("SLIDE", "HEX"))
    parser.add_argument("--hyperlink", nargs=3,
                        metavar=("SLIDE", "TEXT", "URL"))
    parser.add_argument("--enable-slide-number", type=int, metavar="N")
    parser.add_argument("--set-footer", nargs=2, metavar=("SLIDE", "TEXT"))
    parser.add_argument("--set-notes", nargs=2, metavar=("SLIDE", "TEXT"))
    parser.add_argument("--append-notes", nargs=2, metavar=("SLIDE", "TEXT"))
    args = parser.parse_args(argv)

    prs = Presentation(args.pptx)
    report = {"ok": True, "replacements": 0}

    for old, new in args.replace_text:
        report["replacements"] += replace_text(prs, old, new)
    if args.chart_data:
        update_chart(prs, args.chart_data)
        report["chart_updated"] = True
    if args.swap_image:
        swap_image(prs, *args.swap_image)
        report["image_swapped"] = True
    if args.duplicate_slide is not None:
        report["duplicated_to"] = duplicate_slide(prs, args.duplicate_slide)
    if args.set_background:
        set_background(prs.slides[int(args.set_background[0])],
                       args.set_background[1])
        report["background_set"] = True
    if args.hyperlink:
        report["hyperlinked_runs"] = add_hyperlink(prs, *args.hyperlink)
    if args.enable_slide_number is not None:
        report["slide_number_enabled"] = enable_slide_number(
            prs.slides[args.enable_slide_number])
    if args.set_footer:
        set_footer(prs.slides[int(args.set_footer[0])], args.set_footer[1])
        report["footer_set"] = True
    if args.set_notes:
        set_notes(prs.slides[int(args.set_notes[0])], args.set_notes[1])
        report["notes_set"] = True
    if args.append_notes:
        set_notes(prs.slides[int(args.append_notes[0])],
                  args.append_notes[1], append=True)
        report["notes_appended"] = True
    if args.remove_slide is not None:
        remove_slide(prs, args.remove_slide)
        report["slide_removed"] = args.remove_slide
    if args.move_slide:
        move_slide(prs, *args.move_slide)
        report["slide_moved"] = args.move_slide

    out = args.output or args.pptx
    prs.save(out)
    report["output"] = out
    print(json.dumps(report))
    return 0


if __name__ == "__main__":
    sys.exit(main())
